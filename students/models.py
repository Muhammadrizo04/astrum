import time
from django.db import models
import qrcode
from io import BytesIO
from django.core.files import File
from PIL import Image, ImageDraw, ImageFont
from django.conf import settings
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx import Presentation
from django.core.files.base import ContentFile


# class ITEducator(models.Model):
#     first_name = models.CharField(max_length=50)
#     last_name = models.CharField(max_length=50)
#     middle_name = models.CharField(max_length=200)
#     issue_date = models.CharField(max_length=200)
#     expiration_date = models.CharField(max_length=200)
#     create_date = models.DateTimeField(auto_now_add=True)
#     overall_percentage = models.CharField(max_length=30)
#     overall_score = models.CharField(max_length=30)
#     module_1_percentage = models.CharField(max_length=30)
#     module_1_score = models.CharField(max_length=30)
#     module_2_percentage = models.CharField(max_length=30)
#     module_2_score = models.CharField(max_length=30)
#     module_3_percentage = models.CharField(max_length=30)
#     module_3_score = models.CharField(max_length=30)
#     pptx_file = models.FileField(upload_to='pptx_MK', null=True, blank=True)
#     qr_code = models.ImageField(upload_to='malaka_qrcode/', null=True, blank=True)
#     certificate_id = models.CharField(max_length=7, unique=True, null=True, blank=True)
#     certificate_id_numeric = models.IntegerField(unique=True, blank=True)
#     certificate_front = models.FileField(upload_to='malaka_ser_front/', null=True, blank=True)
#     certificate_back = models.FileField(upload_to='malaka_ser_back/', null=True, blank=True)
#     series = models.CharField(max_length=3, default='MK')
#
#     def __str__(self):
#         return self.first_name
#
#     @staticmethod
#     def generate_certificate_id():
#         last_student = ITEducator.objects.order_by('-certificate_id').first()
#         if last_student:
#             last_id_int = int(last_student.certificate_id)
#             if last_id_int >= 1:
#                 new_id_int = last_id_int + 1
#                 new_id_str = str(new_id_int).zfill(7)
#                 return new_id_str
#         return "0000001"
#
#     @staticmethod
#     def generate_certificate_id_numeric():
#         last_student = ITEducator.objects.order_by('-certificate_id_numeric').first()
#         if last_student:
#             last_id_int = last_student.certificate_id_numeric
#             new_id_int = last_id_int + 1
#             return new_id_int
#         return 1
#
#     @staticmethod
#     def add_image(prs, slide_index, image_path, left, top, height):
#         slide = prs.slides[slide_index]
#         image = Image.open(image_path)
#         image = image.convert('RGBA')
#         data = list(image.getdata())
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         new_image = Image.new('RGBA', image.size)
#         new_image.putdata(new_data)
#         image_stream = BytesIO()
#         new_image.save(image_stream, format="PNG")
#         image_stream.seek(0)
#         slide.shapes.add_picture(image_stream, left, top, height=height)
#
#     @staticmethod
#     def add_text(prs, slide_index, left, top, width, height, input_text, font_size, font_color,
#                  font_name='Gilroy', alignment=PP_ALIGN.LEFT):
#         slide = prs.slides[slide_index]
#         text_box = slide.shapes.add_textbox(left, top, width, height)
#         text_frame = text_box.text_frame
#         text = text_frame.add_paragraph()
#         text.text = input_text
#         text.font.size = Pt(font_size)
#         text.font.name = font_name
#         text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
#         text.alignment = alignment
#         for run in text.runs:
#             run.font.bold = True
#
#     def generate_certificate(self):
#         pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/malaka.pptx')
#         prs = Presentation(pptx_template_path)
#
#         # Slide 0
#
#         black_color = (0, 0, 0,)
#         text_color = (0x54, 0x30, 0xCE)
#         text = f"{self.last_name} {self.first_name} {self.middle_name}"
#         series = f"{self.series} {self.certificate_id}"
#         qr_code = os.path.join(settings.MEDIA_ROOT, f'malaka_qrcode/qr_code-{self.certificate_id}.png')
#         qr_code_img_back = os.path.join(settings.MEDIA_ROOT, f'malaka_qrcode/qr_code-{self.certificate_id}.png')
#         issued = f"{self.issue_date}"
#         expiration = f"{self.expiration_date}"
#
#         self.add_image(prs, 0, qr_code, Inches(0.3976377953), Inches(4.7952755906), Inches(1.3))
#         self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
#                       alignment=PP_ALIGN.CENTER)
#         self.add_text(prs, 0, Inches(3.82), Inches(5.64), Inches(1), Inches(1), issued, 12, black_color, )
#         self.add_text(prs, 0, Inches(5.04), Inches(5.64), Inches(1), Inches(1), expiration, 12, black_color, )
#         self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), series, 12, black_color)
#
#         # Slide 1
#
#         module_1_percentage = f"{self.module_1_percentage}%"
#         module_1_score = f"{self.module_1_score} ball/ points"
#         module_2_percentage = f"{self.module_2_percentage}%"
#         module_2_score = f"{self.module_2_score} ball/ points"
#         module_3_percentage = f"{self.module_3_percentage}%"
#         module_3_score = f"{self.module_3_score} ball/ points"
#         overall_percentage = f"{self.overall_percentage}%"
#         overall_score = f"{self.overall_score} ball/ points"
#
#         self.add_text(prs, 1, Inches(0.54), Inches(0.028), Inches(1), Inches(0.8), series, 11, black_color)
#         self.add_text(prs, 1, Inches(1), Inches(0.75), Inches(8), Inches(1), text, 28, text_color,
#                       alignment=PP_ALIGN.CENTER)
#         self.add_text(prs, 1, Inches(2.3622047244), Inches(3.8307086614), Inches(1.7440944882), Inches(0.3346456693),
#                       module_1_percentage, 14, black_color)  # Module 1
#         self.add_text(prs, 1, Inches(2.3622047244), Inches(3.594488189), Inches(1.7440944882), Inches(0.3346456693),
#                       module_1_score, 14, black_color)  # Module 1 Score
#         self.add_text(prs, 1, Inches(2.3622047244), Inches(4.7480314961), Inches(1.7440944882), Inches(0.3346456693),
#                       module_2_percentage, 14, black_color)  # Module 2
#         self.add_text(prs, 1, Inches(2.3622047244), Inches(4.5118110236), Inches(1.7440944882), Inches(0.3346456693),
#                       module_2_score, 14, black_color)  # Module 2 Score
#         self.add_text(prs, 1, Inches(2.3622047244), Inches(5.6692913386), Inches(1.7440944882), Inches(0.3346456693),
#                       module_3_percentage, 14, black_color)  # Module 3
#         self.add_text(prs, 1, Inches(2.3622047244), Inches(5.4330708661), Inches(1.7440944882), Inches(0.3346456693),
#                       module_3_score, 14, black_color)  # Module 3 Score
#         self.add_text(prs, 1, Inches(6.3385826772), Inches(4.6456692913), Inches(1.7440944882), Inches(0.3346456693),
#                       overall_percentage, 14, black_color)  # Overall Percentage
#         self.add_text(prs, 1, Inches(6.3385826772), Inches(3.7007874016), Inches(1.7440944882), Inches(0.3346456693),
#                       overall_score, 14, black_color)  # Overall Score
#         self.add_image(prs, 1, qr_code_img_back, Inches(8.4448818898), Inches(4.8385826772), Inches(1.3))
#
#         # add_image(prs, slide, png, left,top,height)
#
#         # add_text(prs, slide, left, top, width, height, text, size, color)
#
#         pptx_buffer = BytesIO()
#         prs.save(pptx_buffer)
#         pptx_buffer.seek(0)
#
#         return pptx_buffer
#
#     def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
#         background = Image.open(background_path)
#         qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)
#
#         qr_code = qr_code.convert("RGBA")
#         data = qr_code.getdata()
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         qr_code.putdata(new_data)
#
#         background.paste(qr_code, position, qr_code)
#
#         font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
#         draw = ImageDraw.Draw(background)
#         seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)
#
#         black_color = (0, 0, 0)
#         text_color = (0x54, 0x30, 0xCE)
#
#         text = f"{self.last_name} {self.first_name}  {self.middle_name}"
#         draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
#         draw.text((800, 1810), f"{self.series} {self.certificate_id}", fill=black_color, font=seria_font)
#         draw.text((1190, 1810), f"{self.issue_date}", fill=black_color, font=seria_font)
#         draw.text((1540, 1810), f"{self.expiration_date}", fill=black_color, font=seria_font)
#
#         background.save(output_path)
#
#     def overlay_qr_code_back(self, background_path, qr_code_path, output_path, position, qr_size):
#         background = Image.open(background_path)
#         qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)
#
#         qr_code = qr_code.convert("RGBA")
#         data = qr_code.getdata()
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         qr_code.putdata(new_data)
#
#         background.paste(qr_code, position, qr_code)
#
#         font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
#         draw = ImageDraw.Draw(background)
#         seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=50)
#         seria_font_sr = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)
#
#         black_color = (0, 0, 0)
#         text_color = (0x54, 0x30, 0xCE)
#         text = f"{self.last_name} {self.first_name}  {self.middle_name}"
#         draw.text((1500, 440), text, fill=text_color, anchor="ms", font=font)
#         draw.text((190, 128), f"{self.series} {self.certificate_id}", fill=black_color, font=seria_font_sr)
#         draw.text((770, 1190), f"{self.module_1_score} ball/ points", fill=black_color, font=seria_font)
#         draw.text((770, 1260), f"{self.module_1_percentage}%", fill=black_color, font=seria_font)
#         draw.text((770, 1480), f"{self.module_2_score} ball/ points", fill=black_color, font=seria_font)
#         draw.text((770, 1540), f"{self.module_2_percentage}%", fill=black_color, font=seria_font)
#         draw.text((770, 1740), f"{self.module_2_score} ball/ points", fill=black_color, font=seria_font)
#         draw.text((770, 1810), f"{self.module_3_percentage}%", fill=black_color, font=seria_font)
#         draw.text((2000, 1230), f"{self.module_3_score} ball/ points", fill=black_color, font=seria_font)
#         draw.text((2000, 1520), f"{self.overall_percentage}%", fill=black_color, font=seria_font)
#
#         background.save(output_path)
#
#     def save(self, *args, **kwargs):
#         if not self.certificate_id:
#             self.certificate_id = self.generate_certificate_id()
#
#         if not self.certificate_id_numeric:
#             self.certificate_id_numeric = self.generate_certificate_id_numeric()
#
#         # Here, on which link you need to create a qr code, put that link
#         qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/MK{self.certificate_id}")
#         canvas = Image.new('RGB', (380, 380), 'white')
#         canvas.paste(qr_code_img)
#         buffer = BytesIO()
#         canvas.save(buffer, format='PNG')
#         qr_code_file_name = f'qr_code-{self.certificate_id}.png'
#         self.qr_code.save(qr_code_file_name, File(buffer), save=False)
#         canvas.close()
#
#         if not self.certificate_front:
#             background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/malaka-1.png')
#             qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
#                                                     f'malaka_qrcode/qr_code-{self.certificate_id}.png')
#             output_image_path_front = os.path.join(settings.MEDIA_ROOT,
#                                                    f'malaka_ser_front/certificate-{self.certificate_id}.png')
#             self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
#                                        (115, 1435), 390)
#             self.certificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)
#
#         if not self.certificate_back:
#             background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/malaka-2.png')
#             qr_code_image_path_back = os.path.join(settings.MEDIA_ROOT,
#                                                    f'malaka_qrcode/qr_code-{self.certificate_id}.png')
#             output_image_path_back = os.path.join(settings.MEDIA_ROOT,
#                                                   f'malaka_ser_back/certificate-{self.certificate_id}.png')
#             self.overlay_qr_code_back(background_image_path_back, qr_code_image_path_back, output_image_path_back,
#                                       (2535, 1445), 390)
#             self.certificate_back = os.path.relpath(output_image_path_back, settings.MEDIA_ROOT)
#
#             self.certificate_back.save(os.path.basename(output_image_path_back),
#                                        File(open(output_image_path_back, 'rb')))
#
#         if not self.pptx_file:
#             pptx_buffer = self.generate_certificate()
#             self.pptx_file.save(
#                 f'{self.series}-{self.certificate_id}-{self.last_name}-{self.first_name}-{self.middle_name}.pptx',
#                 ContentFile(pptx_buffer.read()), save=False)
#
#         super(ITEducator, self).save(*args, **kwargs)

#
# class InteriorDesign(models.Model):
#     ism = models.CharField(max_length=50)
#     familya = models.CharField(max_length=50)
#     sharif = models.CharField(max_length=50)
#     berilgan_vaqt = models.CharField(max_length=200)
#     create_date = models.DateTimeField(auto_now_add=True)
#     seria = models.CharField(max_length=3, default='3D')
#     sertificate_id = models.CharField(max_length=7, unique=True, blank=True)
#     sertificate_id_numeric = models.IntegerField(unique=True, blank=True)
#     qr_code = models.ImageField(upload_to='3D_qrcode/', blank=True)
#     pptx_file = models.FileField(upload_to='pptx_3D', blank=True)
#     sertificate_front = models.FileField(upload_to='3D_ser_front/', blank=True)
#
#     def __str__(self):
#         return self.ism
#
#     def generate_sertificate_id(self):
#         last_student = InteriorDesign.objects.order_by('-sertificate_id').first()
#         if last_student and last_student.sertificate_id:
#             last_id_int = int(last_student.sertificate_id)
#             new_id_int = last_id_int + 1
#             new_id_str = str(new_id_int).zfill(7)
#             return new_id_str
#         return "0000001"
#
#     def generate_sertificate_id_numeric(self):
#         last_student = InteriorDesign.objects.order_by('-sertificate_id_numeric').first()
#         if last_student:
#             last_id_int = last_student.sertificate_id_numeric
#             new_id_int = last_id_int + 1
#             return new_id_int
#         return 1
#
#     def add_image(self, prs, slide_index, image_path, left, top, height):
#         slide = prs.slides[slide_index]
#         image = Image.open(image_path)
#         image = image.convert('RGBA')
#         data = list(image.getdata())
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         new_image = Image.new('RGBA', image.size)
#         new_image.putdata(new_data)
#         image_stream = BytesIO()
#         new_image.save(image_stream, format="PNG")
#         image_stream.seek(0)
#         slide.shapes.add_picture(image_stream, left, top, height=height)
#
#     def add_text(self, prs, slide_index, left, top, width, height, input_text, font_size, font_color,
#                  font_name='Gilroy', alignment=PP_ALIGN.LEFT):
#         slide = prs.slides[slide_index]
#         text_box = slide.shapes.add_textbox(left, top, width, height)
#         text_frame = text_box.text_frame
#         text = text_frame.add_paragraph()
#         text.text = input_text
#         text.font.size = Pt(font_size)
#         text.font.name = font_name
#         text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
#         text.alignment = alignment
#         for run in text.runs:
#             run.font.bold = True
#
#     def generate_certificate(self):
#         pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/3D.pptx')
#         prs = Presentation(pptx_template_path)
#
#         black_color = (0, 0, 0,)
#         text_color = (0x54, 0x30, 0xCE)
#         text = f"{self.familya} {self.ism}  {self.sharif}"
#         seria = f"{self.seria} {self.sertificate_id}"
#         qr_code = os.path.join(settings.MEDIA_ROOT, f'3D_qrcode/qr_code-{self.sertificate_id}.png')
#         berilgan = f"{self.berilgan_vaqt}"
#
#         self.add_image(prs, 0, qr_code, Inches(0.3858267717), Inches(4.7834645669), Inches(1.3))
#         self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
#                       alignment=PP_ALIGN.CENTER)
#         self.add_text(prs, 0, Inches(4.37), Inches(5.64), Inches(1), Inches(1), berilgan, 11, black_color, )
#         self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), seria, 11, black_color)
#
#         pptx_buffer = BytesIO()
#         prs.save(pptx_buffer)
#         pptx_buffer.seek(0)
#
#         return pptx_buffer
#
#     def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
#         background = Image.open(background_path)
#         qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)
#
#         qr_code = qr_code.convert("RGBA")
#         data = qr_code.getdata()
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         qr_code.putdata(new_data)
#
#         background.paste(qr_code, position, qr_code)
#
#         font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
#         draw = ImageDraw.Draw(background)
#         seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)
#
#         black_color = (0, 0, 0)
#         text_color = (0x54, 0x30, 0xCE)
#
#         text = f"{self.familya} {self.ism}  {self.sharif}"
#         draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
#         draw.text((800, 1810), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
#         draw.text((1350, 1810), f"{self.berilgan_vaqt}", fill=black_color, font=seria_font)
#
#         background.save(output_path)
#
#     def save(self, *args, **kwargs):
#
#         if not self.sertificate_id:
#             self.sertificate_id = self.generate_sertificate_id()
#
#         if not self.sertificate_id_numeric:
#             self.sertificate_id_numeric = self.generate_sertificate_id_numeric()
#
#         qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/3D{self.sertificate_id}")
#         canvas = Image.new('RGB', (380, 380), 'white')
#         canvas.paste(qr_code_img)
#         buffer = BytesIO()
#         canvas.save(buffer, format='PNG')
#         qr_code_file_name = f'qr_code-{self.sertificate_id}.png'
#         self.qr_code.save(qr_code_file_name, File(buffer), save=False)
#         canvas.close()
#
#         if not self.sertificate_front:
#             background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/3D.png')
#             qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
#                                                     f'3D_qrcode/qr_code-{self.sertificate_id}.png')
#             output_image_path_front = os.path.join(settings.MEDIA_ROOT,
#                                                    f'3D_ser_front/certificate-{self.sertificate_id}.png')
#             self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
#                                        (115, 1435), 390)
#             self.sertificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)
#
#         if not self.pptx_file:
#             pptx_buffer = self.generate_certificate()
#             self.pptx_file.save(f'{self.seria}-{self.sertificate_id}-{self.familya}-{self.ism}-{self.sharif}.pptx',
#                                 ContentFile(pptx_buffer.read()), save=False)
#
#         super(InteriorDesign, self).save(*args, **kwargs)

#
# class FullStack(models.Model):
#     ism = models.CharField(max_length=50)
#     familya = models.CharField(max_length=50)
#     sharif = models.CharField(max_length=50)
#     berilgan_vaqt = models.CharField(max_length=200, blank=True)
#     create_date = models.DateTimeField(auto_now_add=True)
#     preseason_web = models.CharField(max_length=6)
#     season_arc = models.CharField(max_length=6)
#     season_fullstack = models.CharField(max_length=6)
#     frontend = models.CharField(blank=True, max_length=6)
#     backend = models.CharField(blank=True, max_length=6)
#     seria = models.CharField(max_length=3, blank=True)
#     sertificate_id = models.CharField(max_length=7, blank=True)
#     sertificate_id_numeric = models.IntegerField(blank=True)
#     frontend_url = models.URLField(blank=True, null=True)
#     backend_url = models.URLField(blank=True, null=True)
#     frontend_qrcode = models.ImageField(upload_to='frontend_qrcode', blank=True)
#     backend_qrcode = models.ImageField(upload_to='backend_qrcode', blank=True)
#     pptx_file = models.FileField(upload_to='pptx_FS', blank=True)
#     qr_code = models.ImageField(upload_to='fullstack_qrcode/', blank=True)
#     sertificate_front = models.FileField(upload_to='fullstack_ser_front/', blank=True)
#     sertificate_back = models.FileField(upload_to='fullstack_ser_back/', blank=True)
#
#     def __str__(self):
#         return self.ism
#
#     def generate_seria(self):
#         if not self.backend_url:
#             return "FD"
#         elif not self.frontend_url:
#             return "BD"
#         else:
#             return "FS"
#
#     def generate_sertificate_id(self):
#         series_prefix = self.seria
#         last_student = FullStack.objects.filter(seria=series_prefix).order_by('-sertificate_id').first()
#         if last_student and last_student.sertificate_id:
#             last_id_int = int(last_student.sertificate_id)
#             new_id_int = last_id_int + 1
#             new_id_str = str(new_id_int).zfill(7)
#             return new_id_str
#         return "0000001"
#
#     def generate_sertificate_id_numeric(self):
#         series_prefix = self.seria
#         last_student = FullStack.objects.filter(seria=series_prefix).order_by('-sertificate_id_numeric').first()
#         if last_student:
#             last_id_int = last_student.sertificate_id_numeric
#             new_id_int = last_id_int + 1
#             return new_id_int
#         return 1
#
#     def add_image(self, prs, slide_index, image_path, left, top, height):
#         slide = prs.slides[slide_index]
#         image = Image.open(image_path)
#         image = image.convert('RGBA')
#         data = list(image.getdata())
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         new_image = Image.new('RGBA', image.size)
#         new_image.putdata(new_data)
#         image_stream = BytesIO()
#         new_image.save(image_stream, format="PNG")
#         image_stream.seek(0)
#         slide.shapes.add_picture(image_stream, left, top, height=height)
#
#     def add_text(self, prs, slide_index, left, top, width, height, input_text, font_size, font_color,
#                  font_name='Gilroy', alignment=PP_ALIGN.LEFT):
#         slide = prs.slides[slide_index]
#         text_box = slide.shapes.add_textbox(left, top, width, height)
#         text_frame = text_box.text_frame
#         text = text_frame.add_paragraph()
#         text.text = input_text
#         text.font.size = Pt(font_size)
#         text.font.name = font_name
#         text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
#         text.alignment = alignment
#         for run in text.runs:
#             run.font.bold = True
#
#     def generate_certificate(self):
#         if not self.frontend_url:
#             pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/backend.pptx')
#         elif not self.backend_url:
#             pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/frontend.pptx')
#         else:
#             pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/Fullstack.pptx')
#         prs = Presentation(pptx_template_path)
#
#         black_color = (0, 0, 0,)
#         text_color = (0x54, 0x30, 0xCE)
#         text = f"{self.familya} {self.ism}  {self.sharif}"
#         seria = f"{self.seria} {self.sertificate_id}"
#         qr_code = os.path.join(settings.MEDIA_ROOT, f'fullstack_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
#         qr_code_img_back = os.path.join(settings.MEDIA_ROOT,
#                                         f'backend_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
#         qr_code_img_front = os.path.join(settings.MEDIA_ROOT,
#                                          f'frontend_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
#         berilgan_vaqt_str = f"{self.berilgan_vaqt}"
#         # Slayd 0
#         self.add_image(prs, 0, qr_code, Inches(0.3976377953), Inches(4.7952755906), Inches(1.3))
#         self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
#                       alignment=PP_ALIGN.CENTER)
#         self.add_text(prs, 0, Inches(4.4173228346), Inches(5.64), Inches(1), Inches(1), berilgan_vaqt_str, 12,
#                       black_color, )
#         self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), seria, 12, black_color)
#
#         preseason = f"{self.preseason_web}%"
#         arc_season = f"{self.season_arc}%"
#         full_stack = f"{self.season_fullstack}%"
#         frontend = f"{self.frontend}%"
#         backend = f"{self.backend}%"
#
#         # Slayd 1
#         self.add_text(prs, 1, Inches(0.54), Inches(0.028), Inches(1), Inches(0.8), seria, 12, black_color)
#         self.add_text(prs, 1, Inches(1), Inches(0.75), Inches(8), Inches(1), text, 28, text_color,
#                       alignment=PP_ALIGN.CENTER)
#         self.add_text(prs, 1, Inches(3.1181102362), Inches(3.4645669291), Inches(1), Inches(1), preseason, 14,
#                       black_color)
#         self.add_text(prs, 1, Inches(5.9212598425), Inches(3.4645669291), Inches(1), Inches(1), arc_season, 14,
#                       black_color)
#         self.add_text(prs, 1, Inches(8.9330708661), Inches(3.4645669291), Inches(1), Inches(1), full_stack, 14,
#                       black_color)
#         self.add_text(prs, 1, Inches(3.82), Inches(4.0433070866), Inches(1), Inches(1), frontend, 14, black_color)
#         self.add_text(prs, 1, Inches(8.90), Inches(4.0433070866), Inches(1), Inches(1), backend, 14, black_color)
#
#         if not self.backend_url:
#             self.add_image(prs, 1, qr_code_img_front, Inches(3.3661417323), Inches(4.842519685), Inches(1.4))
#         if not self.frontend_url:
#             self.add_image(prs, 1, qr_code_img_back, Inches(8.4645669291), Inches(4.842519685), Inches(1.4))
#         elif self.backend_url and self.frontend_url:
#             self.add_image(prs, 1, qr_code_img_back, Inches(8.4645669291), Inches(4.842519685), Inches(1.4))
#             self.add_image(prs, 1, qr_code_img_front, Inches(3.3661417323), Inches(4.842519685), Inches(1.4))
#
#         # add_image(prs, slayd, png, left,top,height)
#
#         # add_text(prs, slayd, left, top, width, height, text, size, color)
#         pptx_buffer = BytesIO()
#         prs.save(pptx_buffer)
#         pptx_buffer.seek(0)
#
#         return pptx_buffer
#
#     def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
#         background = Image.open(background_path)
#         qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)
#
#         qr_code = qr_code.convert("RGBA")
#         data = qr_code.getdata()
#         new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#         qr_code.putdata(new_data)
#
#         background.paste(qr_code, position, qr_code)
#
#         font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
#         draw = ImageDraw.Draw(background)
#         seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)
#
#         black_color = (0, 0, 0)
#         text_color = (0x54, 0x30, 0xCE)
#
#         text = f"{self.familya} {self.ism}  {self.sharif}"
#         draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
#         draw.text((800, 1810), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
#         draw.text((1350, 1810), f"{self.berilgan_vaqt}", fill=black_color, font=seria_font)
#
#         background.save(output_path)
#
#     def overlay_qr_code_back(self, background_path, qr_code_path_front, qr_code_path_back, output_path, position_front,
#                              qr_size_front, position_back, qr_size_back, is_frontend_available=True,
#                              is_backend_available=True):
#         background = Image.open(background_path)
#         draw = ImageDraw.Draw(background)
#
#         if qr_code_path_front is not None and is_frontend_available:
#             frontend_qrcode = Image.open(qr_code_path_front).resize((qr_size_front, qr_size_front), Image.LANCZOS)
#             frontend_qrcode = frontend_qrcode.convert("RGBA")
#             data = frontend_qrcode.getdata()
#             new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#             frontend_qrcode.putdata(new_data)
#             background.paste(frontend_qrcode, position_front, frontend_qrcode)
#
#         if qr_code_path_back is not None and is_backend_available:
#             backend_qrcode = Image.open(qr_code_path_back).resize((qr_size_back, qr_size_back), Image.LANCZOS)
#             backend_qrcode = backend_qrcode.convert("RGBA")
#             data = backend_qrcode.getdata()
#             new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
#             backend_qrcode.putdata(new_data)
#             background.paste(backend_qrcode, position_back, backend_qrcode)
#
#         font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
#         seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, "template/Gilroy-Black.ttf"), size=45)
#
#         black_color = (0, 0, 0)
#         text_color = (0x54, 0x30, 0xCE)
#         text = f"{self.familya} {self.ism}  {self.sharif}"
#         draw.text((1500, 440), text, fill=text_color, anchor="ms", font=font)
#
#         # draw.text((205, 115), f"{self.ism} {self.familya} {self.sharif}", fill=text_color, font=font)
#         draw.text((190, 128), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
#         draw.text((890, 1150), f"{self.preseason_web}%", fill=black_color, font=seria_font)
#         draw.text((1760, 1150), f"{self.season_arc}%", fill=black_color, font=seria_font)
#         draw.text((2700, 1150), f"{self.season_fullstack}%", fill=black_color, font=seria_font)
#         draw.text((1120, 1300), f"{self.frontend}%", fill=black_color, font=seria_font)
#         draw.text((2620, 1300), f"{self.backend}%", fill=black_color, font=seria_font)
#
#         background.save(output_path)
#
    # def save(self, *args, **kwargs):
    #     if not self.seria:
    #         self.seria = self.generate_seria()
    #
    #     if not self.sertificate_id:
    #         self.sertificate_id = self.generate_sertificate_id()
    #
    #     if not self.sertificate_id_numeric:
    #         self.sertificate_id_numeric = self.generate_sertificate_id_numeric()
    #
    #     if self.frontend_url:
    #         qr_code_img_front = qrcode.make(f"{self.frontend_url}")
    #         canvas = Image.new('RGB', (500, 500), 'white')
    #         canvas.paste(qr_code_img_front)
    #         buffer = BytesIO()
    #         canvas.save(buffer, format='PNG')
    #         qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
    #         self.frontend_qrcode.save(qr_code_file_name, File(buffer), save=False)
    #         canvas.close()
    #
    #     if self.backend_url:
    #         qr_code_img_back = qrcode.make(f"{self.backend_url}")
    #         canvas = Image.new('RGB', (500, 500), 'white')
    #         canvas.paste(qr_code_img_back)
    #         buffer = BytesIO()
    #         canvas.save(buffer, format='PNG')
    #         qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
    #         self.backend_qrcode.save(qr_code_file_name, File(buffer), save=False)
    #         canvas.close()
    #
    #     qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/FS{self.sertificate_id}")
    #     canvas = Image.new('RGB', (380, 380), 'white')
    #     canvas.paste(qr_code_img)
    #     buffer = BytesIO()
    #     canvas.save(buffer, format='PNG')
    #     qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
    #     self.qr_code.save(qr_code_file_name, File(buffer), save=False)
    #     canvas.close()
    #
    #     if not self.sertificate_front:
    #         if not self.backend_url:
    #             background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/frontend-1.png')
    #         elif not self.frontend_url:
    #             background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/backend-1.png')
    #         elif self.frontend_url and self.backend_url:
    #             background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/Fullstack-1.png')
    #         qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
    #                                                 f'fullstack_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
    #         output_image_path_front = os.path.join(settings.MEDIA_ROOT,
    #                                                f'fullstack_ser_front/certificate-{self.seria}-{self.sertificate_id}.png')
    #         self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
    #                                    (115, 1435), 390)
    #         self.sertificate_front = output_image_path_front
    #
    #     if not self.sertificate_back:
    #         if not self.backend_url:
    #             background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/frontend-2.png')
    #         if not self.frontend_url:
    #             background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/backend-2.png')
    #         elif self.frontend_url and self.backend_url:
    #             background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/Fullstack-2.png')
    #         qr_code_img_front = None
    #         qr_code_img_back = None
    #         output_image_path_back = None
    #
    #         if self.frontend_url:
    #             qr_code_img_front = os.path.join(settings.MEDIA_ROOT,
    #                                              f'frontend_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
    #         if self.backend_url:
    #             qr_code_img_back = os.path.join(settings.MEDIA_ROOT,
    #                                             f'backend_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
    #
    #         output_image_path_back = os.path.join(settings.MEDIA_ROOT,
    #                                               f'fullstack_ser_back/certificate-{self.seria}-{self.sertificate_id}.png')
    #
    #         self.overlay_qr_code_back(background_image_path_back, qr_code_img_front, qr_code_img_back,
    #                                   output_image_path_back, (1020, 1460), 400, (2550, 1460), 400)
    #
    #         self.sertificate_back = output_image_path_back
#
#         if not self.pptx_file:
#             pptx_buffer = self.generate_certificate()
#             self.pptx_file.save(f'{self.seria}-{self.sertificate_id}-{self.familya}-{self.ism}-{self.sharif}.pptx',
#                                 ContentFile(pptx_buffer.read()), save=False)
#
#         super(FullStack, self).save(*args, **kwargs)
#

class DataSciense(models.Model):
    ism = models.CharField(max_length=50)
    familya = models.CharField(max_length=50)
    sharif = models.CharField(max_length=50)
    berilgan_vaqt = models.CharField(max_length=20, blank=True)
    create_date = models.DateTimeField(auto_now_add=True)
    preseason_data = models.CharField(max_length=6)
    season_arc = models.CharField(max_length=6)
    data_science = models.CharField(blank=True, max_length=6)
    machina_learning = models.CharField(blank=True, max_length=6)
    seria = models.CharField(max_length=3, default='DS')
    sertificate_id = models.CharField(max_length=7, unique=True, blank=True)
    sertificate_id_numeric = models.IntegerField(unique=True, blank=True)
    data_science_url = models.URLField(blank=True, null=True)
    machina_learning_url = models.URLField(blank=True, null=True)
    data_science_qrcode = models.ImageField(upload_to='datascience_qrcode/', blank=True)
    machina_learning_qrcode = models.ImageField(upload_to='machinalr_qrcode/', blank=True)
    qr_code = models.ImageField(upload_to='data_qrcode/', blank=True)
    pptx_file = models.FileField(upload_to='pptx_DT', blank=True)
    sertificate_front = models.FileField(upload_to='data_ser_front/', blank=True)
    sertificate_back = models.FileField(upload_to='data_ser_back/', blank=True)

    def __str__(self):
        return self.ism

    def generate_sertificate_id(self):
        last_student = DataSciense.objects.order_by('-sertificate_id').first()
        if last_student and last_student.sertificate_id:
            last_id_int = int(last_student.sertificate_id)
            new_id_int = last_id_int + 1
            new_id_str = str(new_id_int).zfill(7)
            return new_id_str
        return "0000001"

    def generate_sertificate_id_numeric(self):
        last_student = DataSciense.objects.order_by('-sertificate_id_numeric').first()
        if last_student:
            last_id_int = last_student.sertificate_id_numeric
            new_id_int = last_id_int + 1
            return new_id_int
        return 1

    def add_image(self, prs, slide_index, image_path, left, top, height):
        slide = prs.slides[slide_index]
        image = Image.open(image_path)
        image = image.convert('RGBA')
        data = list(image.getdata())
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        new_image = Image.new('RGBA', image.size)
        new_image.putdata(new_data)
        image_stream = BytesIO()
        new_image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, height=height)

    def add_text(self, prs, slide_index, left, top, width, height, input_text, font_size, font_color,
                 font_name='Gilroy', alignment=PP_ALIGN.LEFT):
        slide = prs.slides[slide_index]
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text = text_frame.add_paragraph()
        text.text = input_text
        text.font.size = Pt(font_size)
        text.font.name = font_name
        text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        text.alignment = alignment
        for run in text.runs:
            run.font.bold = True

    def generate_certificate(self):
        if self.machina_learning_url:
            pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/machine_learning.pptx')
        else:
            pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/data_science.pptx')
        prs = Presentation(pptx_template_path)

        black_color = (0, 0, 0,)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.familya} {self.ism}  {self.sharif}"
        seria = f"{self.seria} {self.sertificate_id}"
        qr_code = os.path.join(settings.MEDIA_ROOT, f'data_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
        qr_code_data = os.path.join(settings.MEDIA_ROOT,
                                    f'datascience_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
        qr_code_machine_learning = os.path.join(settings.MEDIA_ROOT,
                                                f'machinalr_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
        berilgan_vaqt_str = f"{self.berilgan_vaqt}"
        # Slayd 0
        self.add_image(prs, 0, qr_code, Inches(0.3976377953), Inches(4.7952755906), Inches(1.3))
        self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 0, Inches(4.37), Inches(5.64), Inches(1), Inches(1), berilgan_vaqt_str, 11, black_color, )
        self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), seria, 11, black_color)

        prs_data = f"{self.preseason_data}%"
        arc_season = f"{self.season_arc}%"
        data_sc = f"{self.data_science}%"
        machine_lr = f"{self.machina_learning}%"

        # Slayd 1
        self.add_text(prs, 1, Inches(0.54), Inches(0.028), Inches(1), Inches(0.8), seria, 11, black_color)
        self.add_text(prs, 1, Inches(1), Inches(0.75), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 1, Inches(3.82), Inches(3.4645669291), Inches(1), Inches(1), prs_data, 12, black_color)
        self.add_text(prs, 1, Inches(8.90), Inches(3.4645669291), Inches(1), Inches(1), arc_season, 12, black_color)
        self.add_text(prs, 1, Inches(3.82), Inches(4.0433070866), Inches(1), Inches(1), data_sc, 12, black_color)
        self.add_text(prs, 1, Inches(8.90), Inches(4.0433070866), Inches(1), Inches(1), machine_lr, 12, black_color)
        self.add_image(prs, 1, qr_code_data, Inches(3.3661417323), Inches(4.842519685), Inches(1.4))
        if self.machina_learning_url:
            self.add_image(prs, 1, qr_code_machine_learning, Inches(8.4645669291), Inches(4.842519685), Inches(1.4))

        # add_image(prs, slayd, png, left,top,height)

        # add_text(prs, slayd, left, top, width, height, text, size, color)
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)

        text = f"{self.familya} {self.ism}  {self.sharif}"
        draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
        draw.text((800, 1810), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
        draw.text((1350, 1810), f"{self.berilgan_vaqt}", fill=black_color, font=seria_font)

        background.save(output_path)

    def overlay_qr_code_back(self, background_path, qr_code_path_front, qr_code_path_back, output_path, position_front,
                             qr_size_front, position_back, qr_size_back):
        background = Image.open(background_path)
        draw = ImageDraw.Draw(background)

        if qr_code_path_front is not None:
            frontend_qrcode = Image.open(qr_code_path_front).resize((qr_size_front, qr_size_front), Image.LANCZOS)
            frontend_qrcode = frontend_qrcode.convert("RGBA")
            data = frontend_qrcode.getdata()
            new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
            frontend_qrcode.putdata(new_data)
            background.paste(frontend_qrcode, position_front, frontend_qrcode)

        if qr_code_path_back is not None:
            backend_qrcode = Image.open(qr_code_path_back).resize((qr_size_back, qr_size_back), Image.LANCZOS)
            backend_qrcode = backend_qrcode.convert("RGBA")
            data = backend_qrcode.getdata()
            new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
            backend_qrcode.putdata(new_data)
            background.paste(backend_qrcode, position_back, backend_qrcode)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=50)
        seria_font_sr = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, "template/Gilroy-Black.ttf"),
                                           size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.familya} {self.ism}  {self.sharif}"
        draw.text((1500, 440), text, fill=text_color, anchor="ms", font=font)

        # draw.text((205, 115), f"{self.ism} {self.familya} {self.sharif}", fill=text_color, font=font)
        draw.text((190, 128), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font_sr)
        draw.text((1120, 1150), f"{self.preseason_data}%", fill=black_color, font=seria_font)
        draw.text((2700, 1150), f"{self.season_arc}%", fill=black_color, font=seria_font)
        draw.text((1120, 1300), f"{self.data_science}%", fill=black_color, font=seria_font)
        draw.text((2700, 1300), f"{self.machina_learning}%", fill=black_color, font=seria_font)

        background.save(output_path)

    def save(self, *args, **kwargs):
        if not self.sertificate_id:
            self.sertificate_id = self.generate_sertificate_id()

        if not self.sertificate_id_numeric:
            self.sertificate_id_numeric = self.generate_sertificate_id_numeric()
        super(DataSciense, self).save(*args, **kwargs)

        if self.data_science_url:
            qr_code_img_datasc = qrcode.make(f"{self.data_science_url}")
            canvas = Image.new('RGB', (500, 500), 'white')
            canvas.paste(qr_code_img_datasc)
            buffer = BytesIO()
            canvas.save(buffer, format='PNG')
            qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
            self.data_science_qrcode.save(qr_code_file_name, File(buffer), save=False)
            canvas.close()

        if self.machina_learning_url:
            qr_code_img_machinelr = qrcode.make(f"{self.machina_learning_url}")
            canvas = Image.new('RGB', (500, 500), 'white')
            canvas.paste(qr_code_img_machinelr)
            buffer = BytesIO()
            canvas.save(buffer, format='PNG')
            qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
            self.machina_learning_qrcode.save(qr_code_file_name, File(buffer), save=False)
            canvas.close()

        qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
        qr_code_path = os.path.join(settings.MEDIA_ROOT, 'data_qrcode', qr_code_file_name)

        if not os.path.exists(qr_code_path):
            qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/DS{self.sertificate_id}")
            canvas = Image.new('RGB', (380, 380), 'white')
            canvas.paste(qr_code_img)
            buffer = BytesIO()
            canvas.save(buffer, format='PNG')
            # qr_code_file_name = f'qr_code-{self.seria}-{self.sertificate_id}.png'
            self.qr_code.save(qr_code_file_name, File(buffer), save=False)
            canvas.close()

        if not self.sertificate_front:
            background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/data_science-1.png')
            qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                    f'data_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
            output_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                   f'data_ser_front/certificate-{self.seria}-{self.sertificate_id}.png')
            self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
                                       (115, 1435), 390)
            self.sertificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)

        if not self.sertificate_back:
            qr_code_img_front = None
            qr_code_img_back = None
            if self.machina_learning_url:
                background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/machine_learning.png')
            if self.data_science_url:
                background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/data_science-2.png')

            if self.data_science_url:
                qr_code_img_front = os.path.join(settings.MEDIA_ROOT,
                                                 f'datascience_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')
            if self.machina_learning_url:
                qr_code_img_back = os.path.join(settings.MEDIA_ROOT,
                                                f'machinalr_qrcode/qr_code-{self.seria}-{self.sertificate_id}.png')

            if qr_code_img_front or qr_code_img_back:
                output_image_path_back = os.path.join(settings.MEDIA_ROOT,
                                                      f'data_ser_back/certificate-{self.seria}-{self.sertificate_id}.png')

                # Set the coordinates for "Not Available" text

                self.overlay_qr_code_back(background_image_path_back, qr_code_img_front, qr_code_img_back,
                                          output_image_path_back, (1020, 1460), 400, (2550, 1460), 400)

                self.sertificate_back = os.path.relpath(output_image_path_back, settings.MEDIA_ROOT)

        if not self.pptx_file:
            pptx_buffer = self.generate_certificate()
            self.pptx_file.save(f'{self.seria}-{self.sertificate_id}-{self.familya}-{self.ism}-{self.sharif}.pptx',
                                ContentFile(pptx_buffer.read()), save=False)

        super(DataSciense, self).save(*args, **kwargs)


class SoftWare(models.Model):
    ism = models.CharField(max_length=50)
    familya = models.CharField(max_length=50)
    sharif = models.CharField(max_length=50)
    berilgan_vaqt = models.CharField(max_length=20)
    create_date = models.DateTimeField(auto_now_add=True)
    preseason_web = models.CharField(max_length=6)
    season_arc = models.CharField(max_length=6)
    season_arc_2 = models.CharField(blank=True, max_length=6)
    software = models.CharField(blank=True, max_length=6)
    rust = models.CharField(default='0', max_length=6)
    seria = models.CharField(max_length=3, default='SE')
    sertificate_id = models.CharField(max_length=7, unique=True, null=True, blank=True)
    sertificate_id_numeric = models.IntegerField(unique=True, blank=True)
    software_url = models.URLField(blank=True, null=True)
    rust_url = models.URLField(blank=True, null=True)
    software_qrcode = models.ImageField(upload_to='soft_qrcode/', blank=True)
    rust_qrcode = models.ImageField(upload_to='rust_qrcode/', blank=True)
    qr_code = models.ImageField(upload_to='software_qrcode/', blank=True)
    pptx_file = models.FileField(upload_to='pptx_SF', blank=True)
    sertificate_front = models.FileField(upload_to='software_ser_front/', blank=True)
    sertificate_back = models.FileField(upload_to='software_ser_back/', blank=True)

    def __str__(self):
        return self.ism

    def generate_sertificate_id(self):
        last_student = SoftWare.objects.order_by('-sertificate_id').first()
        if last_student and last_student.sertificate_id:
            last_id_int = int(last_student.sertificate_id)
            new_id_int = last_id_int + 1
            new_id_str = str(new_id_int).zfill(7)
            return new_id_str
        return "0000001"

    def generate_sertificate_id_numeric(self):
        last_student = SoftWare.objects.order_by('-sertificate_id_numeric').first()
        if last_student:
            last_id_int = last_student.sertificate_id_numeric
            new_id_int = last_id_int + 1
            return new_id_int
        return 1

    def add_image(self, prs, slide_index, image_path, left, top, height):
        slide = prs.slides[slide_index]
        image = Image.open(image_path)
        image = image.convert('RGBA')
        data = list(image.getdata())
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        new_image = Image.new('RGBA', image.size)
        new_image.putdata(new_data)
        image_stream = BytesIO()
        new_image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, height=height)

    def add_text(self, prs, slide_index, left, top, width, height, input_text, font_size, font_color,
                 font_name='Gilroy', alignment=PP_ALIGN.LEFT):
        slide = prs.slides[slide_index]
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text = text_frame.add_paragraph()
        text.text = input_text
        text.font.size = Pt(font_size)
        text.font.name = font_name
        text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        text.alignment = alignment
        for run in text.runs:
            run.font.bold = True

    def generate_certificate(self):
        pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/software.pptx')
        prs = Presentation(pptx_template_path)

        black_color = (0, 0, 0,)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.familya} {self.ism}  {self.sharif}"
        seria = f"{self.seria} {self.sertificate_id}"
        qr_code = os.path.join(settings.MEDIA_ROOT, f'software_qrcode/qr_code-{self.sertificate_id}.png')
        qr_code_img_front = os.path.join(settings.MEDIA_ROOT, f'soft_qrcode/qr_code-{self.sertificate_id}.png')
        berilgan_vaqt_str = f"{self.berilgan_vaqt}"
        # Slayd 0
        self.add_image(prs, 0, qr_code, Inches(0.3976377953), Inches(4.7952755906), Inches(1.3))
        self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 0, Inches(4.37), Inches(5.64), Inches(1), Inches(1), berilgan_vaqt_str, 11, black_color, )
        self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), seria, 11, black_color)

        prs_web = f"{self.preseason_web}%"
        arc_season = f"{self.season_arc}%"
        arc_season_2 = f"{self.season_arc_2}%"
        software = f"{self.software}%"
        rust = f"{self.rust}%"

        # Slayd 1
        self.add_text(prs, 1, Inches(0.54), Inches(0.028), Inches(1), Inches(0.8), seria, 11, black_color)
        self.add_text(prs, 1, Inches(1), Inches(0.75), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 1, Inches(3.1181102362), Inches(3.4645669291), Inches(1), Inches(1), prs_web, 12,
                      black_color)
        self.add_text(prs, 1, Inches(5.9212598425), Inches(3.4645669291), Inches(1), Inches(1), arc_season, 12,
                      black_color)
        self.add_text(prs, 1, Inches(8.9330708661), Inches(3.4645669291), Inches(1), Inches(1), arc_season_2, 12,
                      black_color)
        self.add_text(prs, 1, Inches(3.82), Inches(4.0433070866), Inches(1), Inches(1), software, 12, black_color)
        self.add_text(prs, 1, Inches(8.90), Inches(4.0433070866), Inches(1), Inches(1), rust, 12, black_color)
        self.add_image(prs, 1, qr_code_img_front, Inches(3.3661417323), Inches(4.842519685), Inches(1.4))

        # add_image(prs, slayd, png, left,top,height)

        # add_text(prs, slayd, left, top, width, height, text, size, color)
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)

        text = f"{self.familya} {self.ism}  {self.sharif}"
        draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
        draw.text((800, 1810), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
        draw.text((1350, 1810), f"{self.berilgan_vaqt}", fill=black_color, font=seria_font)

        background.save(output_path)

    def overlay_qr_code_back(self, background_path, qr_code_path_front, qr_code_path_back, output_path, position_front,
                             qr_size_front, position_back, qr_size_back, is_frontend_available=True,
                             is_backend_available=True):
        background = Image.open(background_path)
        draw = ImageDraw.Draw(background)

        if qr_code_path_front is not None and is_frontend_available:
            frontend_qrcode = Image.open(qr_code_path_front).resize((qr_size_front, qr_size_front), Image.LANCZOS)
            frontend_qrcode = frontend_qrcode.convert("RGBA")
            data = frontend_qrcode.getdata()
            new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
            frontend_qrcode.putdata(new_data)
            background.paste(frontend_qrcode, position_front, frontend_qrcode)

        if qr_code_path_back is not None and is_backend_available:
            backend_qrcode = Image.open(qr_code_path_back).resize((qr_size_back, qr_size_back), Image.LANCZOS)
            backend_qrcode = backend_qrcode.convert("RGBA")
            data = backend_qrcode.getdata()
            new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
            backend_qrcode.putdata(new_data)
            background.paste(backend_qrcode, position_back, backend_qrcode)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=50)
        seria_font_sr = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, "template/Gilroy-Black.ttf"), size=45)

        black_color = (0, 0, 0)
        text_color = (0x53, 0x2F, 0xCE)
        text = f"{self.familya} {self.ism}  {self.sharif}"
        draw.text((1500, 440), text, fill=text_color, anchor="ms", font=font, stroke_width=4)

        # draw.text((205, 115), f"{self.ism} {self.familya} {self.sharif}", fill=text_color, font=font)
        draw.text((190, 128), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font_sr)
        draw.text((890, 1150), f"{self.preseason_web}%", fill=black_color, font=seria_font)
        draw.text((1760, 1150), f"{self.season_arc}%", fill=black_color, font=seria_font)
        draw.text((2700, 1150), f"{self.season_arc_2}%", fill=black_color, font=seria_font)
        draw.text((1120, 1300), f"{self.software}%", fill=black_color, font=seria_font)
        draw.text((2620, 1300), f"{self.rust}%", fill=black_color, font=seria_font)

        background.save(output_path)

    def save(self, *args, **kwargs):

        if not self.sertificate_id:
            self.sertificate_id = self.generate_sertificate_id()

        if not self.sertificate_id_numeric:
            self.sertificate_id_numeric = self.generate_sertificate_id_numeric()

        if self.software_url:
            qr_code_img_front = qrcode.make(f"{self.software_url}")
            canvas = Image.new('RGB', (500, 500), 'white')
            canvas.paste(qr_code_img_front)
            buffer = BytesIO()
            canvas.save(buffer, format='PNG')
            qr_code_file_name = f'qr_code-{self.sertificate_id}.png'
            self.software_qrcode.save(qr_code_file_name, File(buffer), save=False)
            canvas.close()

        qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/SE{self.sertificate_id}")
        canvas = Image.new('RGB', (380, 380), 'white')
        canvas.paste(qr_code_img)
        buffer = BytesIO()
        canvas.save(buffer, format='PNG')
        qr_code_file_name = f'qr_code-{self.sertificate_id}.png'
        self.qr_code.save(qr_code_file_name, File(buffer), save=False)
        canvas.close()

        if not self.sertificate_front:
            background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/software-1.png')
            qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                    f'software_qrcode/qr_code-{self.sertificate_id}.png')
            output_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                   f'software_ser_front/certificate-{self.sertificate_id}.png')
            self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
                                       (115, 1435), 390)
            self.sertificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)

        if not self.sertificate_back:
            background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/software-2.png')
            qr_code_img_front = None
            qr_code_img_back = None
            output_image_path_back = None

            if self.software_url:
                qr_code_img_front = os.path.join(settings.MEDIA_ROOT, f'soft_qrcode/qr_code-{self.sertificate_id}.png')
            if self.rust_url:
                qr_code_img_back = os.path.join(settings.MEDIA_ROOT, f'rust_qrcode/qr_code-{self.sertificate_id}.png')

            if qr_code_img_front or qr_code_img_back:
                output_image_path_back = os.path.join(settings.MEDIA_ROOT,
                                                      f'software_ser_back/certificate-{self.sertificate_id}.png')

                # Set the coordinates for "Not Available" text

                self.overlay_qr_code_back(background_image_path_back, qr_code_img_front, qr_code_img_back,
                                          output_image_path_back, (1020, 1460), 400, (2550, 1460), 400,
                                          is_frontend_available=bool(qr_code_img_front),
                                          is_backend_available=bool(qr_code_img_back))

                self.sertificate_back.save(os.path.basename(output_image_path_back),
                                           File(open(output_image_path_back, 'rb')))

        if not self.pptx_file:
            pptx_buffer = self.generate_certificate()
            self.pptx_file.save(f'{self.seria}-{self.sertificate_id}-{self.familya}-{self.ism}-{self.sharif}.pptx',
                                ContentFile(pptx_buffer.read()), save=False)

        super(SoftWare, self).save(*args, **kwargs)


class Other(models.Model):
    ism = models.CharField(max_length=50)
    familya = models.CharField(max_length=50)
    sharif = models.CharField(max_length=50)
    berilgan_vaqt = models.CharField(max_length=200)
    create_date = models.DateTimeField(auto_now_add=True)
    seria = models.CharField(max_length=3)
    sertificate_id = models.CharField(max_length=7, unique=True, blank=True)
    sertificate_id_numeric = models.IntegerField(unique=True, blank=True)
    pptx_file = models.FileField(upload_to='pptx_other', blank=True)
    sertificate_front = models.FileField(upload_to='other_ser_front/')

    def __str__(self):
        return self.ism


class CyberSecurity(models.Model):
    ism = models.CharField(max_length=50)
    familya = models.CharField(max_length=50)
    sharif = models.CharField(max_length=50, blank=True)
    berilgan_vaqt = models.CharField(max_length=200)
    seria = models.CharField(max_length=3, default='CS')
    sertificate_id = models.CharField(max_length=7, unique=True, blank=True)
    sertificate_id_numeric = models.IntegerField(unique=True, blank=True)
    pptx_file = models.FileField(upload_to='pptx_cyber_security', blank=True)
    sertificate_front = models.FileField(upload_to='cyber_security_front/', blank=True)
    qr_code = models.ImageField(upload_to='cyber_security_qrcode/', blank=True)

    def generate_sertificate_id(self):
        last_student = CyberSecurity.objects.order_by('-sertificate_id').first()
        if last_student and last_student.sertificate_id:
            last_id_int = int(last_student.sertificate_id)
            new_id_int = last_id_int + 1
            new_id_str = str(new_id_int).zfill(7)
            return new_id_str
        return "0000006"

    def generate_sertificate_id_numeric(self):
        last_student = CyberSecurity.objects.order_by('-sertificate_id_numeric').first()
        if last_student:
            last_id_int = last_student.sertificate_id_numeric
            new_id_int = last_id_int + 1
            return new_id_int
        return 6

    def add_image(self, prs, slide_index, image_path, left, top, height):
        slide = prs.slides[slide_index]
        image = Image.open(image_path)
        image = image.convert('RGBA')
        data = list(image.getdata())
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        new_image = Image.new('RGBA', image.size)
        new_image.putdata(new_data)
        image_stream = BytesIO()
        new_image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, height=height)

    def add_text(self, prs, slide_index, left, top, width, height, input_text, font_size, font_color,
                 font_name='Gilroy', alignment=PP_ALIGN.LEFT):
        slide = prs.slides[slide_index]
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text = text_frame.add_paragraph()
        text.text = input_text
        text.font.size = Pt(font_size)
        text.font.name = font_name
        text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        text.alignment = alignment
        for run in text.runs:
            run.font.bold = True

    def generate_certificate(self):
        pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/cyber_security.pptx')
        prs = Presentation(pptx_template_path)

        black_color = (0, 0, 0,)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.familya} {self.ism}  {self.sharif}"
        seria = f"{self.seria} {self.sertificate_id}"
        qr_code = os.path.join(settings.MEDIA_ROOT, f'cyber_security_qrcode/qr_code-{self.sertificate_id}.png')
        berilgan = f"{self.berilgan_vaqt}"

        self.add_image(prs, 0, qr_code, Inches(0.3858267717), Inches(4.7834645669), Inches(1.3))
        self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 0, Inches(4.76), Inches(5.59), Inches(1), Inches(1), berilgan, 11, black_color, )
        self.add_text(prs, 0, Inches(2.8), Inches(5.57), Inches(1), Inches(1), seria, 11, black_color, )

        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)

        text = f"{self.familya} {self.ism}  {self.sharif}"
        draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
        draw.text((860, 1785), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
        draw.text((1450, 1785), f"{self.berilgan_vaqt}", fill=black_color, font=seria_font)

        background.save(output_path)

    def save(self, *args, **kwargs):

        if not self.sertificate_id:
            self.sertificate_id = self.generate_sertificate_id()

        if not self.sertificate_id_numeric:
            self.sertificate_id_numeric = self.generate_sertificate_id_numeric()

        qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/{self.seria}{self.sertificate_id}")
        canvas = Image.new('RGB', (380, 380), 'white')
        canvas.paste(qr_code_img)
        buffer = BytesIO()
        canvas.save(buffer, format='PNG')
        qr_code_file_name = f'qr_code-{self.sertificate_id}.png'
        self.qr_code.save(qr_code_file_name, File(buffer), save=False)
        canvas.close()

        if not self.sertificate_front:
            background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/cyber_security.png')
            qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                    f'cyber_security_qrcode/qr_code-{self.sertificate_id}.png')
            output_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                   f'cyber_security_front/certificate-{self.sertificate_id}.png')
            self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
                                       (115, 1435), 390)
            self.sertificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)

        if not self.pptx_file:
            pptx_buffer = self.generate_certificate()
            self.pptx_file.save(f'{self.seria}-{self.sertificate_id}-{self.familya}-{self.ism}-{self.sharif}.pptx',
                                ContentFile(pptx_buffer.read()), save=False)

        super(CyberSecurity, self).save(*args, **kwargs)
