from django.db import models
import qrcode
from io import BytesIO
from django.core.files import File
from PIL import Image, ImageDraw, ImageFont
from django.conf import settings
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx import Presentation
from django.core.files.base import ContentFile


class ITEducator(models.Model):
    first_name = models.CharField(max_length=50)
    last_name = models.CharField(max_length=50)
    middle_name = models.CharField(max_length=200)
    issue_date = models.CharField(max_length=200)
    expiration_date = models.CharField(max_length=200)
    create_date = models.DateTimeField(auto_now_add=True)
    overall_percentage = models.CharField(max_length=30)
    overall_score = models.CharField(max_length=30)
    module_1_percentage = models.CharField(max_length=30)
    module_1_score = models.CharField(max_length=30)
    module_2_percentage = models.CharField(max_length=30)
    module_2_score = models.CharField(max_length=30)
    module_3_percentage = models.CharField(max_length=30)
    module_3_score = models.CharField(max_length=30)
    pptx_file = models.FileField(upload_to='pptx_MK', null=True, blank=True)
    qr_code = models.ImageField(upload_to='malaka_qrcode/', null=True, blank=True)
    certificate_id = models.CharField(max_length=7, unique=True, null=True, blank=True)
    certificate_id_numeric = models.IntegerField(unique=True, blank=True)
    certificate_front = models.FileField(upload_to='malaka_ser_front/', null=True, blank=True)
    certificate_back = models.FileField(upload_to='malaka_ser_back/', null=True, blank=True)
    series = models.CharField(max_length=3, default='MK')

    def __str__(self):
        return self.first_name

    @staticmethod
    def generate_certificate_id():
        last_student = ITEducator.objects.order_by('-certificate_id').first()
        if last_student:
            last_id_int = int(last_student.certificate_id)
            if last_id_int >= 1:
                new_id_int = last_id_int + 1
                new_id_str = str(new_id_int).zfill(7)
                return new_id_str
        return "0000001"

    @staticmethod
    def generate_certificate_id_numeric():
        last_student = ITEducator.objects.order_by('-certificate_id_numeric').first()
        if last_student:
            last_id_int = last_student.certificate_id_numeric
            new_id_int = last_id_int + 1
            return new_id_int
        return 1

    @staticmethod
    def add_image(prs, slide_index, image_path, left, top, height):
        slide = prs.slides[slide_index]
        image = Image.open(image_path)
        image = image.convert('RGBA')
        data = list(image.getdata())
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        new_image = Image.new('RGBA', image.size)
        new_image.putdata(new_data)
        image_stream = BytesIO()
        new_image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, height=height)

    @staticmethod
    def add_text(prs, slide_index, left, top, width, height, input_text, font_size, font_color,
                 font_name='Gilroy', alignment=PP_ALIGN.LEFT):
        slide = prs.slides[slide_index]
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text = text_frame.add_paragraph()
        text.text = input_text
        text.font.size = Pt(font_size)
        text.font.name = font_name
        text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        text.alignment = alignment
        for run in text.runs:
            run.font.bold = True

    def generate_certificate(self):
        pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/malaka.pptx')
        prs = Presentation(pptx_template_path)

        # Slide 0

        black_color = (0, 0, 0,)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.last_name} {self.first_name} {self.middle_name}"
        series = f"{self.series} {self.certificate_id}"
        qr_code = os.path.join(settings.MEDIA_ROOT, f'malaka_qrcode/qr_code-{self.certificate_id}.png')
        qr_code_img_back = os.path.join(settings.MEDIA_ROOT, f'malaka_qrcode/qr_code-{self.certificate_id}.png')
        issued = f"{self.issue_date}"
        expiration = f"{self.expiration_date}"

        self.add_image(prs, 0, qr_code, Inches(0.3976377953), Inches(4.7952755906), Inches(1.3))
        self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 0, Inches(3.82), Inches(5.64), Inches(1), Inches(1), issued, 12, black_color, )
        self.add_text(prs, 0, Inches(5.04), Inches(5.64), Inches(1), Inches(1), expiration, 12, black_color, )
        self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), series, 12, black_color)

        # Slide 1

        module_1_score = f"{self.module_1_score} ball/ points"
        module_2_score = f"{self.module_2_score} ball/ points"
        module_3_score = f"{self.module_3_score} ball/ points"
        overall_score = f"{self.overall_score} ball/ points"

        self.add_text(prs, 1, Inches(0.54), Inches(0.028), Inches(1), Inches(0.8), series, 11, black_color)
        self.add_text(prs, 1, Inches(1), Inches(0.75), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)

        self.add_text(prs, 1, Inches(2.44), Inches(3.74), Inches(1.7440944882), Inches(0.3346456693),
                      module_1_score, 14, black_color)  # Module 1 Score
        self.add_text(prs, 1, Inches(2.44), Inches(4.61), Inches(1.7440944882), Inches(0.3346456693),
                      module_2_score, 14, black_color)  # Module 2 Score
        self.add_text(prs, 1, Inches(2.44), Inches(5.55), Inches(1.7440944882), Inches(0.3346456693),
                      module_3_score, 14, black_color)  # Module 3 Score
        self.add_text(prs, 1, Inches(6.3385826772), Inches(3.7), Inches(1.7440944882), Inches(0.3346456693),
                      overall_score, 14, black_color)  # Overall Score
        self.add_image(prs, 1, qr_code_img_back, Inches(8.4448818898), Inches(4.8385826772), Inches(1.3))

        # add_image(prs, slide, png, left,top,height)

        # add_text(prs, slide, left, top, width, height, text, size, color)

        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)

        text = f"{self.last_name} {self.first_name}  {self.middle_name}"
        draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
        draw.text((800, 1810), f"{self.series} {self.certificate_id}", fill=black_color, font=seria_font)
        draw.text((1190, 1810), f"{self.issue_date}", fill=black_color, font=seria_font)
        draw.text((1540, 1810), f"{self.expiration_date}", fill=black_color, font=seria_font)

        background.save(output_path)

    def overlay_qr_code_back(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=50)
        seria_font_sr = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.last_name} {self.first_name}  {self.middle_name}"
        draw.text((1500, 440), text, fill=text_color, anchor="ms", font=font)
        draw.text((190, 128), f"{self.series} {self.certificate_id}", fill=black_color, font=seria_font_sr)
        draw.text((770, 1220), f"{self.module_1_score} ball/ points", fill=black_color, font=seria_font)
        draw.text((770, 1510), f"{self.module_2_score} ball/ points", fill=black_color, font=seria_font)
        draw.text((770, 1770), f"{self.module_3_score} ball/ points", fill=black_color, font=seria_font)
        draw.text((1940, 1240), f"{self.overall_score} ball/ points", fill=black_color, font=seria_font)

        background.save(output_path)

    def save(self, *args, **kwargs):
        if not self.certificate_id:
            self.certificate_id = self.generate_certificate_id()

        if not self.certificate_id_numeric:
            self.certificate_id_numeric = self.generate_certificate_id_numeric()

        # Here, on which link you need to create a qr code, put that link
        qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/MK{self.certificate_id}")
        canvas = Image.new('RGB', (380, 380), 'white')
        canvas.paste(qr_code_img)
        buffer = BytesIO()
        canvas.save(buffer, format='PNG')
        qr_code_file_name = f'qr_code-{self.certificate_id}.png'
        self.qr_code.save(qr_code_file_name, File(buffer), save=False)
        canvas.close()

        if not self.certificate_front:
            background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/malaka-1.png')
            qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                    f'malaka_qrcode/qr_code-{self.certificate_id}.png')
            output_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                   f'malaka_ser_front/certificate-{self.certificate_id}.png')
            self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
                                       (115, 1435), 390)
            self.certificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)

        if not self.certificate_back:
            background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/malaka-2.png')
            qr_code_image_path_back = os.path.join(settings.MEDIA_ROOT,
                                                   f'malaka_qrcode/qr_code-{self.certificate_id}.png')
            output_image_path_back = os.path.join(settings.MEDIA_ROOT,
                                                  f'malaka_ser_back/certificate-{self.certificate_id}.png')
            self.overlay_qr_code_back(background_image_path_back, qr_code_image_path_back, output_image_path_back,
                                      (2535, 1445), 390)

            self.certificate_back = output_image_path_back

        if not self.pptx_file:
            pptx_buffer = self.generate_certificate()
            self.pptx_file.save(
                f'{self.series}-{self.certificate_id}-{self.last_name}-{self.first_name}-{self.middle_name}.pptx',
                ContentFile(pptx_buffer.read()), save=False)

        super(ITEducator, self).save(*args, **kwargs)
