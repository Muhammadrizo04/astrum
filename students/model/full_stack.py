from django.db import models
import qrcode
from io import BytesIO
from django.core.files import File
from PIL import Image, ImageDraw, ImageFont
from django.conf import settings
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx import Presentation
from django.core.files.base import ContentFile


class FullStack(models.Model):
    first_name = models.CharField(max_length=50)
    last_name = models.CharField(max_length=50)
    middle_name = models.CharField(max_length=50)
    issue_date = models.CharField(max_length=200, blank=True)
    create_date = models.DateTimeField(auto_now_add=True)
    preseason_web = models.CharField(max_length=6)
    season_arc = models.CharField(max_length=6)
    season_fullstack = models.CharField(max_length=6)
    frontend = models.CharField(blank=True, max_length=6)
    frontend_url = models.URLField(blank=True, null=True)
    backend = models.CharField(blank=True, max_length=6)
    backend_url = models.URLField(blank=True, null=True)
    frontend_qrcode = models.ImageField(upload_to='frontend_qrcode', blank=True)
    backend_qrcode = models.ImageField(upload_to='backend_qrcode', blank=True)
    pptx_file = models.FileField(upload_to='pptx_FS', blank=True)
    qr_code = models.ImageField(upload_to='fullstack_qrcode/', blank=True)
    certificate_front = models.FileField(upload_to='fullstack_ser_front/', blank=True)
    certificate_back = models.FileField(upload_to='fullstack_ser_back/', blank=True)
    series = models.CharField(max_length=3, blank=True)
    certificate_id = models.CharField(max_length=7, blank=True)
    certificate_id_numeric = models.IntegerField(blank=True)

    def __str__(self):
        return self.first_name

    def generate_series(self):
        if not self.backend_url:
            return "FD"
        elif not self.frontend_url:
            return "BD"
        else:
            return "FS"

    def generate_certificate_id(self):
        series_prefix = self.series
        last_student = FullStack.objects.filter(series=series_prefix).order_by('-certificate_id').first()
        if last_student and last_student.certificate_id:
            last_id_int = int(last_student.certificate_id)
            new_id_int = last_id_int + 1
            new_id_str = str(new_id_int).zfill(7)
            return new_id_str
        return "0000001"

    def generate_certificate_id_numeric(self):
        series_prefix = self.series
        last_student = FullStack.objects.filter(series=series_prefix).order_by('-certificate_id_numeric').first()
        if last_student:
            last_id_int = last_student.certificate_id_numeric
            new_id_int = last_id_int + 1
            return new_id_int
        return 1

    @staticmethod
    def add_image(prs, slide_index, image_path, left, top, height):
        slide = prs.slides[slide_index]
        image = Image.open(image_path)
        image = image.convert('RGBA')
        data = list(image.getdata())
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        new_image = Image.new('RGBA', image.size)
        new_image.putdata(new_data)
        image_stream = BytesIO()
        new_image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, height=height)

    @staticmethod
    def add_text(prs, slide_index, left, top, width, height, input_text, font_size, font_color,
                 font_name='Gilroy', alignment=PP_ALIGN.LEFT):
        slide = prs.slides[slide_index]
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text = text_frame.add_paragraph()
        text.text = input_text
        text.font.size = Pt(font_size)
        text.font.name = font_name
        text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        text.alignment = alignment
        for run in text.runs:
            run.font.bold = True

    def generate_certificate(self):
        if not self.frontend_url:
            pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/backend.pptx')
        elif not self.backend_url:
            pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/frontend.pptx')
        else:
            pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/Fullstack.pptx')
        prs = Presentation(pptx_template_path)

        black_color = (0, 0, 0,)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.last_name} {self.first_name}  {self.middle_name}"
        series = f"{self.series} {self.certificate_id}"
        qr_code = os.path.join(settings.MEDIA_ROOT, f'fullstack_qrcode/qr_code-{self.series}-{self.certificate_id}.png')
        qr_code_img_back = os.path.join(settings.MEDIA_ROOT,
                                        f'backend_qrcode/qr_code-{self.series}-{self.certificate_id}.png')
        qr_code_img_front = os.path.join(settings.MEDIA_ROOT,
                                         f'frontend_qrcode/qr_code-{self.series}-{self.certificate_id}.png')
        issue_date_str = f"{self.issue_date}"
        # Slide 0
        self.add_image(prs, 0, qr_code, Inches(0.3976377953), Inches(4.7952755906), Inches(1.3))
        self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 0, Inches(4.4173228346), Inches(5.64), Inches(1), Inches(1), issue_date_str, 12,
                      black_color, )
        self.add_text(prs, 0, Inches(2.55), Inches(5.64), Inches(1), Inches(1), series, 12, black_color)

        preseason = f"{self.preseason_web}%"
        arc_season = f"{self.season_arc}%"
        full_stack = f"{self.season_fullstack}%"
        frontend = f"{self.frontend}%"
        backend = f"{self.backend}%"

        # Slide 1
        self.add_text(prs, 1, Inches(0.54), Inches(0.028), Inches(1), Inches(0.8), series, 12, black_color)
        self.add_text(prs, 1, Inches(1), Inches(0.75), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 1, Inches(3.1181102362), Inches(3.4645669291), Inches(1), Inches(1), preseason, 14,
                      black_color)
        self.add_text(prs, 1, Inches(5.9212598425), Inches(3.4645669291), Inches(1), Inches(1), arc_season, 14,
                      black_color)
        self.add_text(prs, 1, Inches(8.9330708661), Inches(3.4645669291), Inches(1), Inches(1), full_stack, 14,
                      black_color)
        self.add_text(prs, 1, Inches(3.82), Inches(4.0433070866), Inches(1), Inches(1), frontend, 14, black_color)
        self.add_text(prs, 1, Inches(8.90), Inches(4.0433070866), Inches(1), Inches(1), backend, 14, black_color)

        if not self.backend_url:
            self.add_image(prs, 1, qr_code_img_front, Inches(3.3661417323), Inches(4.842519685), Inches(1.4))
        if not self.frontend_url:
            self.add_image(prs, 1, qr_code_img_back, Inches(8.4645669291), Inches(4.842519685), Inches(1.4))
        elif self.backend_url and self.frontend_url:
            self.add_image(prs, 1, qr_code_img_back, Inches(8.4645669291), Inches(4.842519685), Inches(1.4))
            self.add_image(prs, 1, qr_code_img_front, Inches(3.3661417323), Inches(4.842519685), Inches(1.4))

        # add_image(prs, slide, png, left,top,height)

        # add_text(prs, slide, left, top, width, height, text, size, color)
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        series_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)

        text = f"{self.last_name} {self.first_name}  {self.middle_name}"
        draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
        draw.text((800, 1810), f"{self.series} {self.certificate_id}", fill=black_color, font=series_font)
        draw.text((1350, 1810), f"{self.issue_date}", fill=black_color, font=series_font)

        background.save(output_path)

    def overlay_qr_code_back(self, background_path, qr_code_path_front, qr_code_path_back, output_path, position_front,
                             qr_size_front, position_back, qr_size_back, is_frontend_available=True,
                             is_backend_available=True):
        background = Image.open(background_path)
        draw = ImageDraw.Draw(background)

        if qr_code_path_front is not None and is_frontend_available:
            frontend_qrcode = Image.open(qr_code_path_front).resize((qr_size_front, qr_size_front), Image.LANCZOS)
            frontend_qrcode = frontend_qrcode.convert("RGBA")
            data = frontend_qrcode.getdata()
            new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
            frontend_qrcode.putdata(new_data)
            background.paste(frontend_qrcode, position_front, frontend_qrcode)

        if qr_code_path_back is not None and is_backend_available:
            backend_qrcode = Image.open(qr_code_path_back).resize((qr_size_back, qr_size_back), Image.LANCZOS)
            backend_qrcode = backend_qrcode.convert("RGBA")
            data = backend_qrcode.getdata()
            new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
            backend_qrcode.putdata(new_data)
            background.paste(backend_qrcode, position_back, backend_qrcode)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        series_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, "template/Gilroy-Black.ttf"), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.last_name} {self.first_name}  {self.middle_name}"
        draw.text((1500, 440), text, fill=text_color, anchor="ms", font=font)

        # draw.text((205, 115), f"{self.first_name} {self.last_name} {self.middle_name}", fill=text_color, font=font)
        draw.text((190, 128), f"{self.series} {self.certificate_id}", fill=black_color, font=series_font)
        draw.text((890, 1150), f"{self.preseason_web}%", fill=black_color, font=series_font)
        draw.text((1760, 1150), f"{self.season_arc}%", fill=black_color, font=series_font)
        draw.text((2700, 1150), f"{self.season_fullstack}%", fill=black_color, font=series_font)
        draw.text((1120, 1300), f"{self.frontend}%", fill=black_color, font=series_font)
        draw.text((2620, 1300), f"{self.backend}%", fill=black_color, font=series_font)

        background.save(output_path)

    def save(self, *args, **kwargs):
        background_image_path_front = None
        background_image_path_back = None
        if not self.series:
            self.series = self.generate_series()

        if not self.certificate_id:
            self.certificate_id = self.generate_certificate_id()

        if not self.certificate_id_numeric:
            self.certificate_id_numeric = self.generate_certificate_id_numeric()

        if self.frontend_url:
            qr_code_img_front = qrcode.make(f"{self.frontend_url}")
            canvas = Image.new('RGB', (500, 500), 'white')
            canvas.paste(qr_code_img_front)
            buffer = BytesIO()
            canvas.save(buffer, format='PNG')
            qr_code_file_name = f'qr_code-{self.series}-{self.certificate_id}.png'
            self.frontend_qrcode.save(qr_code_file_name, File(buffer), save=False)
            canvas.close()

        if self.backend_url:
            qr_code_img_back = qrcode.make(f"{self.backend_url}")
            canvas = Image.new('RGB', (500, 500), 'white')
            canvas.paste(qr_code_img_back)
            buffer = BytesIO()
            canvas.save(buffer, format='PNG')
            qr_code_file_name = f'qr_code-{self.series}-{self.certificate_id}.png'
            self.backend_qrcode.save(qr_code_file_name, File(buffer), save=False)
            canvas.close()

        qr_code_img = qrcode.make(f"172.20.31.57:8000/student/{self.series}{self.certificate_id}")
        canvas = Image.new('RGB', (380, 380), 'white')
        canvas.paste(qr_code_img)
        buffer = BytesIO()
        canvas.save(buffer, format='PNG')
        qr_code_file_name = f'qr_code   -{self.series}-{self.certificate_id}.png'
        self.qr_code.save(qr_code_file_name, File(buffer), save=False)
        canvas.close()

        if not self.certificate_front:
            if not self.backend_url:
                background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/frontend-1.png')
            elif not self.frontend_url:
                background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/backend-1.png')
            elif self.frontend_url and self.backend_url:
                background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/Fullstack-1.png')
            qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                    f'fullstack_qrcode/qr_code-{self.series}-{self.certificate_id}.png')
            output_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                   f'fullstack_ser_front/certificate-{self.series}-{self.certificate_id}.png')
            self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
                                       (115, 1435), 390)
            self.certificate_front = output_image_path_front

        if not self.certificate_back:
            if not self.backend_url:
                background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/frontend-2.png')
            if not self.frontend_url:
                background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/backend-2.png')
            elif self.frontend_url and self.backend_url:
                background_image_path_back = os.path.join(settings.MEDIA_ROOT, 'template/Fullstack-2.png')
            qr_code_img_front = None
            qr_code_img_back = None

            if self.frontend_url:
                qr_code_img_front = os.path.join(settings.MEDIA_ROOT,
                                                 f'frontend_qrcode/qr_code-{self.series}-{self.certificate_id}.png')
            if self.backend_url:
                qr_code_img_back = os.path.join(settings.MEDIA_ROOT,
                                                f'backend_qrcode/qr_code-{self.series}-{self.certificate_id}.png')

            output_image_path_back = os.path.join(settings.MEDIA_ROOT,
                                                  f'fullstack_ser_back/certificate-{self.series}-{self.certificate_id}.png')

            self.overlay_qr_code_back(background_image_path_back, qr_code_img_front, qr_code_img_back,
                                      output_image_path_back, (1020, 1460), 400, (2550, 1460), 400)

            self.certificate_back = output_image_path_back

        if not self.pptx_file:
            pptx_buffer = self.generate_certificate()
            self.pptx_file.save(
                f'{self.series}-{self.certificate_id}-{self.last_name}-{self.first_name}-{self.middle_name}.pptx',
                ContentFile(pptx_buffer.read()), save=False)

        super(FullStack, self).save(*args, **kwargs)
