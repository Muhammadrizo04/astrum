from django.db import models
import qrcode
from io import BytesIO
from django.core.files import File
from PIL import Image, ImageDraw, ImageFont
from django.conf import settings
import os
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from pptx import Presentation
from django.core.files.base import ContentFile

class NetworkAdmin(models.Model):
    ism = models.CharField(max_length=50)
    familya = models.CharField(max_length=50)
    sharif = models.CharField(max_length=50)
    berilgan_vaqt = models.CharField(max_length=12)
    amal_qilish = models.CharField(max_length=12)
    seria = models.CharField(max_length=3, default='NA')
    sertificate_id = models.CharField(max_length=7, unique=True, blank=True)
    sertificate_id_numeric = models.IntegerField(unique=True, blank=True)
    pptx_file = models.FileField(upload_to='pptx_networkadmin', blank=True)
    qr_code = models.ImageField(upload_to='ccna_qrcode/', blank=True)
    sertificate_front = models.FileField(upload_to='networkadmin_ser_front/', blank=True)

    def __str__(self):
        return self.ism


    def generate_sertificate_id(self):
        last_student = NetworkAdmin.objects.order_by('-sertificate_id').first()
        if last_student and last_student.sertificate_id:
            last_id_int = int(last_student.sertificate_id)
            new_id_int = last_id_int + 1
            new_id_str = str(new_id_int).zfill(7)
            return new_id_str
        return "0000001"

    def generate_sertificate_id_numeric(self):
        last_student = NetworkAdmin.objects.order_by('-sertificate_id_numeric').first()
        if last_student:
            last_id_int = last_student.sertificate_id_numeric
            new_id_int = last_id_int + 1
            return new_id_int
        return 6

    def add_image(self, prs, slide_index, image_path, left, top, height):
        slide = prs.slides[slide_index]
        image = Image.open(image_path)
        image = image.convert('RGBA')
        data = list(image.getdata())
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        new_image = Image.new('RGBA', image.size)
        new_image.putdata(new_data)
        image_stream = BytesIO()
        new_image.save(image_stream, format="PNG")
        image_stream.seek(0)
        slide.shapes.add_picture(image_stream, left, top, height=height)

    def add_text(self, prs, slide_index, left, top, width, height, input_text, font_size, font_color,
                 font_name='Gilroy', alignment=PP_ALIGN.LEFT):
        slide = prs.slides[slide_index]
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text = text_frame.add_paragraph()
        text.text = input_text
        text.font.size = Pt(font_size)
        text.font.name = font_name
        text.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        text.alignment = alignment
        for run in text.runs:
            run.font.bold = True

    def generate_certificate(self):
        pptx_template_path = os.path.join(settings.MEDIA_ROOT, 'template/ccna.pptx')
        prs = Presentation(pptx_template_path)

        black_color = (0, 0, 0,)
        text_color = (0x54, 0x30, 0xCE)
        text = f"{self.familya} {self.ism}  {self.sharif}"
        seria = f"{self.seria} {self.sertificate_id}"
        qr_code = os.path.join(settings.MEDIA_ROOT, f'ccna_qrcode/qr_code-{self.sertificate_id}.png')
        berilgan = f"{self.berilgan_vaqt}"
        amal_qilish = f"{self.amal_qilish}"

        self.add_image(prs, 0, qr_code, Inches(0.3858267717), Inches(4.7834645669), Inches(1.3))
        self.add_text(prs, 0, Inches(1), Inches(2.55), Inches(8), Inches(1), text, 28, text_color,
                      alignment=PP_ALIGN.CENTER)
        self.add_text(prs, 0, Inches(5.04), Inches(5.65), Inches(1), Inches(1), amal_qilish, 12, black_color, )
        self.add_text(prs, 0, Inches(3.86), Inches(5.65), Inches(1), Inches(1), berilgan, 12, black_color, )
        self.add_text(prs, 0, Inches(2.6), Inches(5.66), Inches(1), Inches(1), seria, 11, black_color, )

        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)

        return pptx_buffer

    def overlay_qr_code_front(self, background_path, qr_code_path, output_path, position, qr_size):
        background = Image.open(background_path)
        qr_code = Image.open(qr_code_path).resize((qr_size, qr_size), Image.LANCZOS)

        qr_code = qr_code.convert("RGBA")
        data = qr_code.getdata()
        new_data = [(255, 255, 255, 0) if item[:3] == (255, 255, 255) else item for item in data]
        qr_code.putdata(new_data)

        background.paste(qr_code, position, qr_code)

        font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=100)
        draw = ImageDraw.Draw(background)
        seria_font = ImageFont.truetype(os.path.join(settings.MEDIA_ROOT, 'template/Gilroy-Black.ttf'), size=45)

        black_color = (0, 0, 0)
        text_color = (0x54, 0x30, 0xCE)

        text = f"{self.familya} {self.ism}  {self.sharif}"
        draw.text((1500, 980), text, fill=text_color, anchor="ms", font=font)
        draw.text((800, 1810), f"{self.seria} {self.sertificate_id}", fill=black_color, font=seria_font)
        draw.text((1190, 1810), f"{self.berilgan_vaqt}", fill=black_color, font=seria_font)
        draw.text((1540, 1810), f"{self.amal_qilish }", fill=black_color, font=seria_font)

        background.save(output_path)

    def save(self, *args, **kwargs):

        if not self.sertificate_id:
            self.sertificate_id = self.generate_sertificate_id()

        if not self.sertificate_id_numeric:
            self.sertificate_id_numeric = self.generate_sertificate_id_numeric()

        qr_code_img = qrcode.make(f"https://certificate.astrum.uz/student/NA{self.sertificate_id}")
        canvas = Image.new('RGB', (380, 380), 'white')
        canvas.paste(qr_code_img)
        buffer = BytesIO()
        canvas.save(buffer, format='PNG')
        qr_code_file_name = f'qr_code-{self.sertificate_id}.png'
        self.qr_code.save(qr_code_file_name, File(buffer), save=False)
        canvas.close()

        if not self.sertificate_front:
            background_image_path_front = os.path.join(settings.MEDIA_ROOT, 'template/ccna.png')
            qr_code_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                    f'ccna_qrcode/qr_code-{self.sertificate_id}.png')
            output_image_path_front = os.path.join(settings.MEDIA_ROOT,
                                                   f'networkadmin_ser_front/certificate-{self.sertificate_id}.png')
            self.overlay_qr_code_front(background_image_path_front, qr_code_image_path_front, output_image_path_front,
                                       (115, 1435), 390)
            self.sertificate_front = os.path.relpath(output_image_path_front, settings.MEDIA_ROOT)

        if not self.pptx_file:
            pptx_buffer = self.generate_certificate()
            self.pptx_file.save(f'{self.seria}-{self.sertificate_id}-{self.familya}-{self.ism}-{self.sharif}.pptx',
                                ContentFile(pptx_buffer.read()), save=False)

        super(NetworkAdmin, self).save(*args, **kwargs)